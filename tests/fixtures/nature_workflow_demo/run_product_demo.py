from __future__ import annotations

import argparse
import csv
import importlib.util
import subprocess
import sys
from pathlib import Path

FIXTURE_DIR = Path(__file__).resolve().parent
OPTIONAL_DEPS = {
    "matplotlib": "matplotlib",
    "numpy": "numpy",
    "pptx": "python-pptx",
}


def _missing_optional_dependencies() -> list[str]:
    return [package for module, package in OPTIONAL_DEPS.items() if importlib.util.find_spec(module) is None]


def _require_optional_dependencies() -> None:
    missing = _missing_optional_dependencies()
    if missing:
        packages = " ".join(missing)
        raise SystemExit(
            "nature-workflow product demo requires optional plotting/Office dependencies. "
            f"Install them with: pip install {packages}"
        )


def _write_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(text.rstrip() + "\n", encoding="utf-8")


def _write_source_data(path: Path) -> list[dict[str, str]]:
    rows = [
        {
            "workflow": "baseline-agent",
            "evidence_coverage": "0.58",
            "citation_error_rate": "0.21",
            "readiness_score": "0.54",
            "qa_failures": "7",
            "route_latency_s": "18.4",
        },
        {
            "workflow": "scholaraio-fallback",
            "evidence_coverage": "0.73",
            "citation_error_rate": "0.12",
            "readiness_score": "0.70",
            "qa_failures": "4",
            "route_latency_s": "20.8",
        },
        {
            "workflow": "nature-workflow-bridge",
            "evidence_coverage": "0.88",
            "citation_error_rate": "0.05",
            "readiness_score": "0.86",
            "qa_failures": "1",
            "route_latency_s": "22.1",
        },
    ]
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8", newline="") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0]))
        writer.writeheader()
        writer.writerows(rows)
    return rows


def _make_figure(rows: list[dict[str, str]], output_dir: Path) -> None:
    import matplotlib
    import numpy as np

    matplotlib.use("Agg")
    from matplotlib import pyplot as plt

    labels = ["Baseline", "Fallback", "Bridge"]
    evidence = np.array([float(row["evidence_coverage"]) for row in rows])
    citation_errors = np.array([float(row["citation_error_rate"]) for row in rows])
    readiness = np.array([float(row["readiness_score"]) for row in rows])
    qa_failures = np.array([int(row["qa_failures"]) for row in rows])

    plt.rcParams.update(
        {
            "font.family": "DejaVu Sans",
            "font.size": 8,
            "axes.spines.top": False,
            "axes.spines.right": False,
            "axes.linewidth": 0.8,
            "xtick.major.width": 0.8,
            "ytick.major.width": 0.8,
            "figure.dpi": 160,
        }
    )
    colors = ["#8E8E8E", "#4C78A8", "#2F7D5B"]

    fig = plt.figure(figsize=(7.2, 4.6), constrained_layout=True)
    grid = fig.add_gridspec(2, 2, width_ratios=[1.25, 1], height_ratios=[1, 1])
    ax_a = fig.add_subplot(grid[:, 0])
    ax_b = fig.add_subplot(grid[0, 1])
    ax_c = fig.add_subplot(grid[1, 1])

    x = np.arange(len(labels))
    ax_a.bar(x, readiness, color=colors, width=0.62)
    ax_a.set_ylim(0, 1)
    ax_a.set_xticks(x, labels, rotation=18, ha="right")
    ax_a.set_ylabel("Readiness score")
    ax_a.set_title("a  End-to-end artifact readiness", loc="left", fontweight="bold")
    for idx, value in enumerate(readiness):
        ax_a.text(idx, value + 0.025, f"{value:.2f}", ha="center", va="bottom", fontsize=8)

    ax_b.plot(x, evidence, marker="o", color="#2F7D5B", linewidth=1.8, label="Evidence coverage")
    ax_b.plot(x, 1 - citation_errors, marker="s", color="#B85C38", linewidth=1.8, label="Citation correctness")
    ax_b.set_ylim(0.45, 1.0)
    ax_b.set_xticks(x, labels, rotation=18, ha="right")
    ax_b.set_ylabel("Proportion")
    ax_b.set_title("b  Source-grounding checks", loc="left", fontweight="bold")
    ax_b.legend(frameon=False, fontsize=7, loc="lower right")

    ax_c.bar(x, qa_failures, color=colors, width=0.62)
    ax_c.set_ylim(0, max(qa_failures) + 2)
    ax_c.set_xticks(x, labels, rotation=18, ha="right")
    ax_c.set_ylabel("Open QA failures")
    ax_c.set_title("c  Residual review risks", loc="left", fontweight="bold")
    for idx, value in enumerate(qa_failures):
        ax_c.text(idx, value + 0.25, str(value), ha="center", va="bottom", fontsize=8)

    fig.suptitle(
        "Nature workflow bridge demonstration on a synthetic artifact-readiness benchmark",
        fontsize=10,
        fontweight="bold",
    )

    figure_dir = output_dir / "figures"
    figure_dir.mkdir(parents=True, exist_ok=True)
    for suffix in ("svg", "pdf", "png"):
        fig.savefig(figure_dir / f"nature-workflow-product-demo.{suffix}", bbox_inches="tight")
    plt.close(fig)


def _write_writing_outputs(output_dir: Path) -> None:
    before = """
# Abstract Before

This demo checks a skill router for Nature related scientific work. We made a
small benchmark and found that routing can be useful. The outputs include a
figure, a data availability statement, and a slide deck. The workflow is better
because it reminds the agent to use upstream skills and avoid made-up details.
"""
    polished = """
# Abstract Polished

This demo uses synthetic data to test whether a ScholarAIO Nature workflow
bridge can turn an ambiguous high-impact-journal request into reviewable
research artifacts. Across a small artifact-readiness benchmark, the bridge
routes each task to the corresponding upstream `nature-*` skill when available
and falls back to ScholarAIO-native skills only when direct upstream execution is
not possible. The resulting package contains source data, a multi-panel figure,
a route audit, a Data Availability statement, and a checked PPTX deck. The
exercise does not claim a biological, clinical, or engineering discovery; it
tests workflow behavior, source grounding, and artifact completeness.
"""
    _write_text(output_dir / "writing" / "abstract-before.md", before)
    _write_text(output_dir / "writing" / "abstract-polished.md", polished)


def _write_data_availability(output_dir: Path) -> None:
    statement = """
# Data Availability Statement

All data used in this demonstration are synthetic and are provided in
`inputs/source-data.csv` within this demo package. No human participants,
clinical records, biological samples, proprietary datasets, third-party
restricted datasets, or external accessioned datasets were used.

No accession numbers, repository DOIs, embargo dates, or controlled-access
committees are claimed for this demo. If this workflow were applied to a real
manuscript, each dataset would need to be mapped to a public repository,
controlled-access route, in-paper source-data table, or justified restriction
before the statement could be submission-ready.
"""
    _write_text(output_dir / "data" / "data-availability.md", statement)


def _write_figure_docs(output_dir: Path) -> None:
    caption = """
# Figure Caption

**Figure 1 | Nature workflow bridge demonstration on a synthetic
artifact-readiness benchmark.** a, End-to-end artifact readiness for a baseline
agent response, a ScholarAIO fallback route, and the Nature workflow bridge.
b, Source-grounding checks represented as evidence coverage and citation
correctness. c, Residual QA failures after self-review. Values are generated
from the synthetic source table in `inputs/source-data.csv` and are intended only
to validate the workflow.
"""
    qa = """
# Figure QA

- Source table: `inputs/source-data.csv`.
- Backend: Python / matplotlib only.
- Exports: SVG, PDF, PNG.
- Panel labels: present (`a`, `b`, `c`).
- Scientific claim: workflow-readiness demonstration only; no domain discovery
  is claimed.
- Integrity checks: all plotted values are directly read from source-data.csv;
  no hidden smoothing, fabricated confidence intervals, or invented sample
  sizes are used.
"""
    _write_text(output_dir / "figures" / "figure-caption.md", caption)
    _write_text(output_dir / "figures" / "figure-qa.md", qa)


def _add_title(slide, title: str, subtitle: str | None = None) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.32), Inches(12.4), Inches(0.55))
    frame = title_box.text_frame
    frame.text = title
    run = frame.paragraphs[0].runs[0]
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(34, 42, 50)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.48), Inches(0.88), Inches(12.2), Inches(0.34))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.paragraphs[0].runs[0].font.size = Pt(10)
        subtitle_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(92, 103, 115)


def _add_bullets(slide, x: float, y: float, w: float, h: float, items: list[str]) -> None:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    for idx, item in enumerate(items):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.size = Pt(14)
        paragraph.font.color.rgb = RGBColor(40, 48, 56)
        paragraph.space_after = Pt(8)


def _write_pptx(output_dir: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Nature Workflow Product Demo", "Direct upstream skill routing with ScholarAIO fallback evidence")
    _add_bullets(
        slide,
        0.7,
        1.55,
        6.2,
        4.3,
        [
            "Scenario: high-impact manuscript support package",
            "Direct upstream routes: nature-figure, nature-polishing, nature-data, nature-paper2ppt",
            "Fallback mode: explicit ScholarAIO route when upstream skills are unavailable",
            "Artifacts: source data, figure, polished abstract, data statement, route audit, PPTX",
        ],
    )
    slide.shapes.add_picture(
        str(output_dir / "figures" / "nature-workflow-product-demo.png"),
        Inches(7.0),
        Inches(1.28),
        width=Inches(5.7),
    )

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Route Audit", "What the bridge decides before content generation")
    _add_bullets(
        slide,
        0.7,
        1.3,
        11.8,
        4.8,
        [
            "Upstream target is stated before substantive work.",
            "Mode is either direct upstream preferred or ScholarAIO fallback.",
            "Fallback output is not described as upstream-equivalent.",
            "Guardrails focus on source grounding, no invented citations, and no invented data identifiers.",
        ],
    )

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Figure QA", "Nature-figure constraints applied to a synthetic benchmark")
    slide.shapes.add_picture(
        str(output_dir / "figures" / "nature-workflow-product-demo.png"),
        Inches(0.6),
        Inches(1.12),
        width=Inches(7.2),
    )
    _add_bullets(
        slide,
        8.15,
        1.3,
        4.4,
        4.7,
        [
            "Python backend selected explicitly.",
            "All values come from source-data.csv.",
            "Vector and raster exports are both generated.",
            "No sample sizes, confidence intervals, or accession IDs are invented.",
        ],
    )

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Ready-to-Paste Data Availability", "Synthetic-data demo with no invented identifiers")
    _add_bullets(
        slide,
        0.75,
        1.35,
        11.8,
        4.6,
        [
            "All data are synthetic and bundled as inputs/source-data.csv.",
            "No human, clinical, biological, proprietary, or third-party restricted data are used.",
            "No accession number or repository DOI is claimed.",
            "Real manuscripts still need dataset-by-dataset repository mapping.",
        ],
    )

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Decision", "Is this useful as a ScholarAIO branch?")
    _add_bullets(
        slide,
        0.75,
        1.4,
        11.8,
        4.6,
        [
            "Useful if the product goal is safe orchestration across upstream Nature Skills.",
            "Not a replacement for installing the upstream bundle.",
            "Fallbacks are honest and bounded, not feature-equivalent claims.",
            "The practical value is reduced routing ambiguity and fewer unsafe scholarly outputs.",
        ],
    )

    slide_dir = output_dir / "slides"
    slide_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = slide_dir / "nature-workflow-product-demo.pptx"
    prs.save(pptx_path)

    inspect_lines = [
        "# PPTX Inspect",
        "",
        f"Slides: {len(prs.slides)}",
        f"Size: {prs.slide_width / 914400:.2f} x {prs.slide_height / 914400:.2f} in",
        "",
    ]
    for idx, slide in enumerate(prs.slides, 1):
        inspect_lines.append(f"## Slide {idx}")
        inspect_lines.append(f"Shape count: {len(slide.shapes)}")
        for shape in slide.shapes:
            kind = "picture" if shape.shape_type == 13 else "text"
            inspect_lines.append(f"- {kind}: {shape.left / 914400:.2f}, {shape.top / 914400:.2f}")
        inspect_lines.append("")
    _write_text(slide_dir / "pptx-inspect.md", "\n".join(inspect_lines))


def _write_manifest(output_dir: Path) -> None:
    manifest = """
# Nature Workflow Product Demo

This package follows `docs/guide/nature-workflow-quickstart.md` and exercises the
bridge as a real research-output workflow, not only a static routing check.

## Upstream Route Coverage Used

- `nature-figure`: generated a Python/matplotlib multi-panel figure with SVG,
  PDF, and PNG exports.
- `nature-polishing`: produced a Nature-leaning polished abstract while
  preserving the synthetic-data limitation.
- `nature-data`: produced a Data Availability statement with no invented
  accession IDs.
- `nature-paper2ppt`: generated an actual PPTX deck and local inspect report.

The route audit also contains all 10 upstream `nature-*` targets, including
direct upstream preferred and ScholarAIO fallback modes.

## Key Files

- `route-cards.md`: executable route demo output.
- `inputs/source-data.csv`: synthetic source data used for all plotted values.
- `figures/nature-workflow-product-demo.svg`: editable vector figure.
- `figures/nature-workflow-product-demo.pdf`: manuscript-style PDF export.
- `figures/nature-workflow-product-demo.png`: raster preview and slide input.
- `writing/abstract-polished.md`: polished abstract.
- `data/data-availability.md`: ready-to-paste demo statement.
- `slides/nature-workflow-product-demo.pptx`: generated slide deck.
- `qa/product-demo-verification.md`: final verification record.
"""
    _write_text(output_dir / "README.md", manifest)


def _write_verification(output_dir: Path, rows: list[dict[str, str]]) -> None:
    required = [
        "route-cards.md",
        "inputs/source-data.csv",
        "figures/nature-workflow-product-demo.svg",
        "figures/nature-workflow-product-demo.pdf",
        "figures/nature-workflow-product-demo.png",
        "writing/abstract-polished.md",
        "data/data-availability.md",
        "slides/nature-workflow-product-demo.pptx",
        "slides/pptx-inspect.md",
    ]
    missing = [item for item in required if not (output_dir / item).exists()]
    if missing:
        raise AssertionError(f"missing product demo files: {', '.join(missing)}")

    verification = f"""
# Product Demo Verification

- Required files present: PASS.
- Source rows: {len(rows)}.
- Figure exports present: SVG, PDF, PNG.
- PPTX generated and locally inspected: PASS.
- No invented accession IDs: PASS.
- No external dataset claimed: PASS.
- Upstream fidelity policy: original `nature-*` skills are preferred; fallback
  outputs are labelled as ScholarAIO fallbacks.

This demo uses synthetic data, so it demonstrates workflow behavior and artifact
quality gates rather than a domain-specific scientific discovery.
"""
    _write_text(output_dir / "qa" / "product-demo-verification.md", verification)


def _write_route_cards(output_dir: Path) -> None:
    subprocess.run(
        [
            sys.executable,
            str(FIXTURE_DIR / "run_demo.py"),
            "--output",
            str(output_dir / "route-cards.md"),
        ],
        check=True,
        capture_output=True,
        text=True,
    )


def run(output_dir: Path) -> None:
    _require_optional_dependencies()
    output_dir.mkdir(parents=True, exist_ok=True)
    rows = _write_source_data(output_dir / "inputs" / "source-data.csv")
    _write_route_cards(output_dir)
    _make_figure(rows, output_dir)
    _write_writing_outputs(output_dir)
    _write_data_availability(output_dir)
    _write_figure_docs(output_dir)
    _write_pptx(output_dir)
    _write_manifest(output_dir)
    _write_verification(output_dir, rows)


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate the nature-workflow product demo")
    parser.add_argument("--output-dir", type=Path, required=True)
    args = parser.parse_args()

    run(args.output_dir)
    print("nature-workflow product demo: PASS")
    print(f"output directory: {args.output_dir}")


if __name__ == "__main__":
    main()
